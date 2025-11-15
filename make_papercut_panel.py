#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Papercut/stencil panel builder from a single image (3x4 A4).

Pipeline (no neural nets):
- Load image, apply optional resolution normalization, convert to grayscale,
  autocontrast, optional Gaussian blur.
- Two-threshold segmentation:
    * threshold-bg: отделить фон от "бумаги".
    * threshold-detail: отделить бумагу от внутренних прорезей/теней.
- Из них строится бинарная маска "бумаги": белые области = то, что остаётся из вытинанки.
- Небольшая морфологическая обработка (closing) для замыкания тонких контуров.
- Собирается панель "белое на сером фоне".
- Панель режется на 3x4 A4 @ DPI (без полей в самой панели).
- На основе тайлов собираются PPTX/PDF с реальными полями (margin-mm).

Вывод:
- PNG исходного размера: сохраняется рядом с оригиналом.
- Панель PNG: <output>_panel.png (если включён --debug-panel).
- PPTX, PDF или оба: <output>_3x4_A4.pptx / .pdf.

Зависимости: pillow, reportlab, python-pptx
"""

from __future__ import annotations

import argparse
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from typing import List, Tuple

import logging
from collections import deque
import sys

from PIL import Image, ImageChops, ImageDraw, ImageFilter, ImageOps, ImageStat

try:  # Pillow>=9.1
    Resampling = Image.Resampling  # type: ignore[attr-defined]
except AttributeError:  # pragma: no cover - fallback for older Pillow
    class _Resampling:
        NEAREST = Image.NEAREST
        BILINEAR = Image.BILINEAR
        BICUBIC = Image.BICUBIC
        LANCZOS = getattr(Image, "LANCZOS", Image.BICUBIC)

    Resampling = _Resampling()  # type: ignore[misc]

# Опционально импортируем pdf/pptx только когда надо
try:
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfgen import canvas
except Exception:  # pragma: no cover
    A4 = None
    canvas = None

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except Exception:  # pragma: no cover
    Presentation = None
    Inches = None
    Pt = None


# --------------------------- геометрия A4 ----------------------------

MM_PER_INCH = 25.4


def mm_to_px(mm: float, dpi: int) -> int:
    return int(round(mm * dpi / MM_PER_INCH))


def px_to_mm(px: int, dpi: int) -> float:
    if dpi <= 0 or px <= 0:
        return 0.0
    return (px / float(dpi)) * MM_PER_INCH


def a4_size_px(dpi: int) -> Tuple[int, int]:
    """A4 portrait in pixels at given DPI."""
    # A4: 210 x 297 mm (портрет)
    w_mm, h_mm = 210.0, 297.0
    return mm_to_px(w_mm, dpi), mm_to_px(h_mm, dpi)


@dataclass
class PanelLayout:
    dpi: int
    cols: int
    rows: int
    margin_mm: float

    @property
    def page_size_px(self) -> Tuple[int, int]:
        return a4_size_px(self.dpi)

    @property
    def tile_size_px(self) -> Tuple[int, int]:
        pw, ph = self.page_size_px
        margin_px = mm_to_px(self.margin_mm, self.dpi)
        inner_w = pw - 2 * margin_px
        inner_h = ph - 2 * margin_px
        if inner_w <= 0 or inner_h <= 0:
            raise ValueError("Margins are too large for A4 at this DPI")
        return inner_w, inner_h

    @property
    def panel_size_px(self) -> Tuple[int, int]:
        tw, th = self.tile_size_px
        return self.cols * tw, self.rows * th


# --------------------------- B/W подготовка ----------------------------


@dataclass
class NormalizationParams:
    target_dpi: int
    upscale_factor: float
    blur_radius: float


@dataclass
class BWParams:
    threshold_bg: int
    threshold_detail: int
    blur: float
    dilate_px: int
    antialias_radius: float
    detail_join_px: int


@dataclass
class MaskResult:
    paper: Image.Image
    filled: Image.Image
    holes: Image.Image
    outline: Image.Image


logger = logging.getLogger("papercut_panel")


def normalize_source_image(
    img: Image.Image, params: NormalizationParams
) -> Tuple[Image.Image, float, float]:
    """Return a smoothed copy of the source image ready for mask extraction."""

    if params.upscale_factor <= 1.0 and params.blur_radius <= 0:
        return img.copy(), 1.0, 0.0

    upscale = max(1.0, float(params.upscale_factor))

    dpi = img.info.get("dpi") if isinstance(img.info.get("dpi"), tuple) else None
    if dpi and params.target_dpi > 0:
        src_dpi = max(dpi[0], dpi[1], 1)
        dpi_scale = params.target_dpi / float(src_dpi)
        upscale = max(upscale, dpi_scale)

    w, h = img.size
    if w == 0 or h == 0:
        return img.copy(), 1.0, 0.0

    target_w = max(1, int(round(w * upscale)))
    target_h = max(1, int(round(h * upscale)))

    if target_w == w and target_h == h:
        working = img.copy()
    else:
        working = img.resize((target_w, target_h), resample=Resampling.LANCZOS)

    applied_blur = 0.0
    if params.blur_radius > 0:
        # Reduce smoothing when we already have sufficient resolution.
        effective_scale = max(upscale, 1.0)
        blur_radius = params.blur_radius / effective_scale
        if dpi and params.target_dpi > 0 and dpi[0] >= params.target_dpi and dpi[1] >= params.target_dpi:
            blur_radius = 0.0
        if blur_radius >= 0.05:
            working = working.filter(ImageFilter.GaussianBlur(radius=blur_radius))
            applied_blur = blur_radius

    if (target_w, target_h) != (w, h):
        working = working.resize((w, h), resample=Resampling.LANCZOS)

    return working, upscale, applied_blur


def prepare_grayscale(gray: Image.Image, params: BWParams) -> Image.Image:
    """Autocontrast + optional blur for robust thresholding."""
    result = ImageOps.autocontrast(gray)
    if params.blur > 0:
        result = result.filter(ImageFilter.GaussianBlur(radius=params.blur))
    return result


def _apply_closing(mask: Image.Image, radius: int) -> Image.Image:
    if radius <= 0:
        return mask
    size = max(3, 2 * radius + 1)  # ensure odd kernel
    mask_l = mask.convert("L")
    mask_l = mask_l.filter(ImageFilter.MaxFilter(size=size))
    mask_l = mask_l.filter(ImageFilter.MinFilter(size=size))
    return mask_l.point(lambda v: 255 if v >= 128 else 0, mode="1")


def extract_outline(mask: Image.Image, radius: int = 1) -> Image.Image:
    """Return a thin outline for the provided binary mask."""

    if radius < 1:
        radius = 1

    mask_bin = mask.convert("1")
    mask_l = mask_bin.convert("L")

    size = max(3, 2 * radius + 1)
    eroded = mask_l.filter(ImageFilter.MinFilter(size=size))
    eroded_bin = eroded.point(lambda v: 255 if v >= 128 else 0, mode="1")

    outline = ImageChops.logical_and(mask_bin, ImageChops.invert(eroded_bin))
    return outline


def _fill_closed_regions_python(blocked_binary: Image.Image) -> Image.Image:
    """Slow but compatible BFS-based fill used as a fallback."""

    w, h = blocked_binary.size
    if w == 0 or h == 0:
        return Image.new("1", (w, h), 0)

    blocked_px = blocked_binary.load()
    outside = Image.new("1", (w, h), 0)
    outside_px = outside.load()
    q = deque()

    def enqueue(x: int, y: int) -> None:
        if 0 <= x < w and 0 <= y < h and blocked_px[x, y] == 0 and outside_px[x, y] == 0:
            outside_px[x, y] = 255
            q.append((x, y))

    for x in range(w):
        enqueue(x, 0)
        enqueue(x, h - 1)
    for y in range(h):
        enqueue(0, y)
        enqueue(w - 1, y)

    neighbours = ((-1, 0), (1, 0), (0, -1), (0, 1))

    while q:
        x, y = q.popleft()
        for dx, dy in neighbours:
            enqueue(x + dx, y + dy)

    solid = ImageChops.invert(outside)
    interior = ImageChops.logical_and(solid, ImageChops.invert(blocked_binary))
    return interior


def fill_closed_regions(outline: Image.Image, *, dilation_radius: int) -> Image.Image:
    """Fill areas enclosed by the outline mask (returns interior only).

    Uses Pillow's flood-fill (C implementation) for performance instead of a
    Python-level BFS so very large panels finish quickly.
    """

    if outline.mode != "1":
        outline = outline.convert("1")

    if dilation_radius > 0:
        size = max(3, 2 * dilation_radius + 1)
        blocked_l = outline.convert("L").filter(ImageFilter.MaxFilter(size=size))
    else:
        blocked_l = outline.convert("L")

    blocked_binary = blocked_l.point(lambda v: 255 if v > 0 else 0, mode="1")

    w, h = blocked_binary.size
    if w == 0 or h == 0:
        return Image.new("1", (w, h), 0)

    flood_img = blocked_binary.convert("L")

    try:
        ImageDraw.floodfill(flood_img, (0, 0), 128, thresh=0)
        outside = flood_img.point(lambda v: 255 if v == 128 else 0, mode="1")
    except AttributeError:  # pragma: no cover - very old Pillow
        logger.debug("ImageDraw.floodfill unavailable, using Python fallback")
        return _fill_closed_regions_python(blocked_binary)

    solid = ImageChops.invert(outside)
    interior = ImageChops.logical_and(solid, ImageChops.invert(blocked_binary))
    return interior


def soften_mask(mask: Image.Image, radius: float) -> Image.Image:
    """Return a crisp alpha mask with optional pre-blur smoothing."""

    mask_l = mask.convert("L")
    if radius <= 0:
        return mask_l.point(lambda v: 255 if v >= 128 else 0, mode="L")

    blurred = mask_l.filter(ImageFilter.GaussianBlur(radius=radius))
    return blurred.point(lambda v: 255 if v >= 128 else 0, mode="L")


def smooth_binary_edges(mask: Image.Image, radius: float) -> Image.Image:
    """Binary smoothing helper used for morphological joins."""

    if radius <= 0:
        return mask.convert("1")
    mask_l = mask.convert("L")
    blurred = mask_l.filter(ImageFilter.GaussianBlur(radius=radius))
    return blurred.point(lambda v: 255 if v >= 128 else 0, mode="1")


def connect_detail_gaps(mask: Image.Image, join_radius: int, *, smooth_radius: float) -> Image.Image:
    """Bridge dotted interior lines into continuous detail regions."""

    base = mask.convert("1")
    result = base

    if join_radius > 0:
        closed = _apply_closing(base, join_radius)
        result = ImageChops.logical_or(result, closed)

    if smooth_radius > 0:
        result = smooth_binary_edges(result, smooth_radius)

    return result


def detect_relative_dark_regions(
    gray_processed: Image.Image,
    paper_candidate: Image.Image,
    t_bg: int,
    t_detail: int,
) -> Image.Image:
    """Detect locally dark regions that should become interior cut-outs."""

    if paper_candidate.mode != "1":
        paper_candidate = paper_candidate.convert("1")

    if paper_candidate.getbbox() is None:
        return Image.new("1", gray_processed.size, 0)

    local_mean = gray_processed.filter(ImageFilter.BoxBlur(1))
    local_min = gray_processed.filter(ImageFilter.MinFilter(size=3))
    local_drop = ImageChops.subtract(local_mean, local_min)

    relative_margin = max(6, (t_bg - t_detail) // 3)
    rel_candidates = local_drop.point(lambda d: 255 if d >= relative_margin else 0, mode="1")

    near_bg = gray_processed.point(lambda g: 255 if g < t_bg - 2 else 0, mode="1")
    rel_candidates = ImageChops.logical_and(rel_candidates, near_bg)
    rel_candidates = ImageChops.logical_and(rel_candidates, paper_candidate)

    edge_map = gray_processed.filter(ImageFilter.FIND_EDGES)
    edge_thresh = max(20, min(96, (t_bg - t_detail) * 2 + 16))
    edge_mask = edge_map.point(lambda v: 255 if v >= edge_thresh else 0, mode="1")
    edge_mask = ImageChops.logical_and(edge_mask, paper_candidate)
    if edge_mask.getbbox() is not None:
        edge_mask = _apply_closing(edge_mask, 1)
        edge_regions = fill_closed_regions(edge_mask, dilation_radius=1)
        rel_candidates = ImageChops.logical_and(rel_candidates, edge_regions)

    rel_closed = _apply_closing(rel_candidates, 1)
    return rel_closed


def scale_radius(base_radius: int, scale_factor: float, *, exponent: float = 0.5, clamp: int | None = None) -> int:
    """Scale morphology radii sub-linearly to preserve fine details when upscaling."""

    if base_radius <= 0:
        return 0

    effective_scale = max(scale_factor, 1.0)
    scaled = int(round(base_radius * (effective_scale ** exponent)))

    if clamp is not None:
        scaled = min(scaled, clamp)

    return max(1, scaled)


def build_masks_from_gray(gray_processed: Image.Image, params: BWParams) -> MaskResult:
    t_bg = max(0, min(255, params.threshold_bg))
    t_detail = max(0, min(t_bg, params.threshold_detail))

    outline_raw = gray_processed.point(lambda g: 255 if g < t_bg else 0, mode="1")
    outline_closed = _apply_closing(outline_raw, params.dilate_px)
    combined_outline = ImageChops.logical_or(outline_raw, outline_closed)

    fill_radius = params.dilate_px if params.dilate_px > 0 else 0
    filled = fill_closed_regions(outline_closed, dilation_radius=fill_radius)
    paper_candidate = filled

    holes_raw = gray_processed.point(lambda g: 255 if g < t_detail else 0, mode="1")
    holes_seed = ImageChops.logical_and(holes_raw, paper_candidate)

    # Relative darkness detector: treat locally dark pockets as potential cut-outs
    holes_rel = detect_relative_dark_regions(gray_processed, paper_candidate, t_bg, t_detail)
    if holes_rel.getbbox() is not None:
        holes_seed = ImageChops.logical_or(holes_seed, holes_rel)

    hole_close_radius = params.dilate_px if params.dilate_px > 0 else 0
    holes_refined = _apply_closing(holes_seed, hole_close_radius)

    detail_outline = extract_outline(holes_refined, radius=1)
    outline_bridge = max(params.dilate_px, params.detail_join_px)
    detail_outline_closed = _apply_closing(detail_outline, outline_bridge)
    detail_filled = fill_closed_regions(detail_outline_closed, dilation_radius=outline_bridge)

    holes_combined = ImageChops.logical_or(holes_refined, detail_filled)
    join_radius = max(0, params.detail_join_px)
    join_smooth = max(0.0, join_radius * 0.35)
    holes_joined = connect_detail_gaps(holes_combined, join_radius, smooth_radius=join_smooth)
    holes_limited = ImageChops.logical_and(holes_joined, filled)

    paper = ImageChops.logical_and(paper_candidate, ImageChops.invert(holes_limited))

    return MaskResult(
        paper=paper,
        filled=filled,
        holes=holes_limited,
        outline=combined_outline,
    )


def build_paper_mask(gray: Image.Image, params: BWParams, *, invert: bool = False) -> Image.Image:
    """
    Сборка бинарной маски "бумаги" из градаций серого.

    gray: L-image 0..255
    возвращает image mode '1' (0=фон/дырка, 255=бумага).
    """
    prepared = prepare_grayscale(gray, params)
    if invert:
        prepared = ImageOps.invert(prepared)
    masks = build_masks_from_gray(prepared, params)
    softened = soften_mask(masks.paper, params.antialias_radius)
    return softened.point(lambda v: 255 if v >= 128 else 0, mode="1")


def _compute_mask_ratio(mask: Image.Image) -> float:
    mask_l = mask.convert("L")
    stat = ImageStat.Stat(mask_l)
    total_pixels = mask.size[0] * mask.size[1]
    if total_pixels == 0:
        return 0.0
    white_pixels = stat.sum[0] / 255.0
    return white_pixels / float(total_pixels)


def _border_and_center_means(gray: Image.Image, border_fraction: float = 0.12) -> Tuple[float, float]:
    w, h = gray.size
    if w == 0 or h == 0:
        return 0.0, 0.0

    bw = max(1, int(round(w * border_fraction)))
    bh = max(1, int(round(h * border_fraction)))

    top_box = (0, 0, w, bh)
    bottom_box = (0, h - bh, w, h)
    left_box = (0, bh, bw, h - bh)
    right_box = (w - bw, bh, w, h - bh)

    center_box = (bw, bh, w - bw, h - bh)

    def _box_mean(box: Tuple[int, int, int, int]) -> Tuple[float, int]:
        x0, y0, x1, y1 = box
        if x0 >= x1 or y0 >= y1:
            return 0.0, 0
        region = gray.crop(box)
        stat = ImageStat.Stat(region)
        return stat.mean[0], (x1 - x0) * (y1 - y0)

    border_parts = [top_box, bottom_box, left_box, right_box]
    border_sum = 0.0
    border_area = 0
    for part in border_parts:
        mean, area = _box_mean(part)
        border_sum += mean * area
        border_area += area

    if border_area == 0:
        border_mean = 0.0
    else:
        border_mean = border_sum / border_area

    center_mean, center_area = _box_mean(center_box)
    if center_area == 0:
        center_mean = ImageStat.Stat(gray).mean[0]

    return border_mean, center_mean


def _should_flip_auto(
    gray_processed: Image.Image,
    mask_result: MaskResult,
    invert_mode: str,
    *,
    tolerance: float = 5.0,
) -> bool:
    mode = invert_mode.lower()
    if mode == "flip":
        return True
    if mode == "keep":
        return False

    paper_ratio = _compute_mask_ratio(mask_result.paper)
    if paper_ratio <= 0.5:
        # Выглядит как фигура меньшего размера, инверсия не нужна.
        return False

    border_mean, center_mean = _border_and_center_means(gray_processed)
    return border_mean + tolerance < center_mean


def build_white_on_gray_panel(
    img: Image.Image,
    layout: PanelLayout,
    bw_params: BWParams,
    bg_rgb: Tuple[int, int, int],
    fit_mode: str = "fit",
    shift_x_mm: float = 0.0,
    shift_y_mm: float = 0.0,
    invert_mode: str = "auto",
) -> Tuple[Image.Image, Image.Image, bool]:
    """
    Основной этап: делаем панель "белое на сером".

    Returns tuple: (panel image, original-size panel, was_auto_flipped).
    """
    logger.info("Preparing grayscale source and preview masks")

    src = img.convert("RGB")

    panel_w, panel_h = layout.panel_size_px

    gray = src.convert("L")
    processed_gray = prepare_grayscale(gray, bw_params)
    preview_masks = build_masks_from_gray(processed_gray, bw_params)

    should_flip = _should_flip_auto(processed_gray, preview_masks, invert_mode)
    if should_flip:
        logger.info("Auto inversion enabled for silhouette")
        processed_gray = ImageOps.invert(processed_gray)
        base_masks = build_masks_from_gray(processed_gray, bw_params)
    else:
        base_masks = preview_masks

    logger.info("Rendering original-size panel preview")
    original_masks = base_masks
    original_mask_soft = soften_mask(original_masks.paper, bw_params.antialias_radius)
    original_panel = Image.new("RGB", processed_gray.size, color=bg_rgb)
    original_tmp = Image.new("RGB", processed_gray.size, color=(255, 255, 255))
    original_panel.paste(original_tmp, (0, 0), original_mask_soft)

    src_w, src_h = processed_gray.size
    if src_w == 0 or src_h == 0:
        raise ValueError("Input image has zero dimension")

    scale_x = panel_w / src_w
    scale_y = panel_h / src_h

    if fit_mode == "stretch":
        scale_used = max(scale_x, scale_y)
        new_w, new_h = panel_w, panel_h
    else:
        if fit_mode == "fill":
            scale_used = max(scale_x, scale_y)
        else:  # default fit
            scale_used = min(scale_x, scale_y)
        new_w = int(round(src_w * scale_used))
        new_h = int(round(src_h * scale_used))

    logger.info("Resizing silhouette to %dx%d (scale %.3f)", new_w, new_h, scale_used)

    scale_factor = max(scale_used, 1.0)
    max_panel_dilate = (
        max(1, mm_to_px(2.0, layout.dpi)) if bw_params.dilate_px > 0 else 0
    )
    max_panel_join = (
        max(1, mm_to_px(1.4, layout.dpi)) if bw_params.detail_join_px > 0 else 0
    )
    scaled_dilate = scale_radius(
        bw_params.dilate_px,
        scale_factor,
        exponent=0.45,
        clamp=max_panel_dilate if max_panel_dilate > 0 else None,
    )
    scaled_join = scale_radius(
        bw_params.detail_join_px,
        scale_factor,
        exponent=0.35,
        clamp=max_panel_join if max_panel_join > 0 else None,
    )

    scaled_params = BWParams(
        threshold_bg=bw_params.threshold_bg,
        threshold_detail=bw_params.threshold_detail,
        blur=bw_params.blur,
        dilate_px=scaled_dilate,
        antialias_radius=bw_params.antialias_radius,
        detail_join_px=scaled_join,
    )

    gray_resized = processed_gray.resize((new_w, new_h), resample=Resampling.LANCZOS)
    dilate_mm = px_to_mm(scaled_dilate, layout.dpi)
    join_mm = px_to_mm(scaled_join, layout.dpi)
    logger.info(
        "Building final masks at export scale (dilate=%d px / %.2f mm, detail_join=%d px / %.2f mm)",
        scaled_dilate,
        dilate_mm,
        scaled_join,
        join_mm,
    )
    final_masks = build_masks_from_gray(gray_resized, scaled_params)

    antialias = max(0.0, bw_params.antialias_radius * (scale_factor ** 0.5))
    antialias_cap = mm_to_px(0.45, layout.dpi)
    if antialias_cap > 0:
        antialias = min(antialias, float(antialias_cap))
    logger.info("Applying antialias radius %.2f", antialias)
    mask_l = soften_mask(final_masks.paper, antialias)

    panel = Image.new("RGB", (panel_w, panel_h), color=bg_rgb)

    shift_x_px = mm_to_px(shift_x_mm, layout.dpi)
    shift_y_px = mm_to_px(shift_y_mm, layout.dpi)

    off_x = (panel_w - new_w) // 2 + shift_x_px
    off_y = (panel_h - new_h) // 2 + shift_y_px

    tmp = Image.new("RGB", (new_w, new_h), color=(255, 255, 255))
    panel.paste(tmp, (off_x, off_y), mask_l)

    return panel, original_panel, should_flip


# --------------------------- нарезка панели ----------------------------


def slice_panel_into_tiles(panel: Image.Image, layout: PanelLayout) -> List[Image.Image]:
    tw, th = layout.tile_size_px
    tiles: List[Image.Image] = []
    for row in range(layout.rows):
        for col in range(layout.cols):
            left = col * tw
            upper = row * th
            box = (left, upper, left + tw, upper + th)
            tiles.append(panel.crop(box))
    return tiles


# --------------------------- вывод PPTX / PDF ----------------------------


def save_pdf_from_tiles(
    tiles: List[Image.Image],
    out_pdf: Path,
    layout: PanelLayout,
    bg_rgb: Tuple[int, int, int],
) -> None:
    if canvas is None or A4 is None:
        raise RuntimeError("reportlab is not installed")

    pw_px, ph_px = layout.page_size_px
    pw_mm, ph_mm = 210.0, 297.0
    c = canvas.Canvas(str(out_pdf), pagesize=A4)

    margin_px = mm_to_px(layout.margin_mm, layout.dpi)
    inner_w_px, inner_h_px = layout.tile_size_px

    for tile in tiles:
        tmp_path = out_pdf.with_suffix(".tmp_tile.png")
        tile.save(tmp_path, format="PNG")

        tile_w_in = inner_w_px / layout.dpi
        tile_h_in = inner_h_px / layout.dpi
        tile_w_pt = tile_w_in * 72.0
        tile_h_pt = tile_h_in * 72.0

        margin_in = (margin_px / layout.dpi)
        margin_pt = margin_in * 72.0

        x = margin_pt
        y = (A4[1] - tile_h_pt) - margin_pt

        c.drawImage(str(tmp_path), x, y, width=tile_w_pt, height=tile_h_pt)
        c.showPage()

        tmp_path.unlink(missing_ok=True)

    c.save()


def save_pptx_from_tiles(
    tiles: List[Image.Image],
    out_pptx: Path,
    layout: PanelLayout,
    bg_rgb: Tuple[int, int, int],
) -> None:
    if Presentation is None or Inches is None:
        raise RuntimeError("python-pptx is not installed")

    prs = Presentation()
    pw_in = 210.0 / MM_PER_INCH
    ph_in = 297.0 / MM_PER_INCH
    prs.slide_width = int(pw_in * 914400)
    prs.slide_height = int(ph_in * 914400)

    margin_px = mm_to_px(layout.margin_mm, layout.dpi)
    inner_w_px, inner_h_px = layout.tile_size_px

    image_streams: List[BytesIO] = []

    for tile in tiles:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        buffer = BytesIO()
        tile.save(buffer, format="PNG")
        buffer.seek(0)
        image_streams.append(buffer)

        tile_w_in = inner_w_px / layout.dpi
        tile_h_in = inner_h_px / layout.dpi
        margin_in = margin_px / layout.dpi

        left = Inches(margin_in)
        top = Inches(margin_in)
        width = Inches(tile_w_in)
        height = Inches(tile_h_in)

        slide.shapes.add_picture(buffer, left, top, width=width, height=height)

    prs.save(out_pptx)


# --------------------------- CLI ----------------------------


def parse_args(argv=None) -> argparse.Namespace:
    p = argparse.ArgumentParser(description="Papercut/stencil panel builder (3x4 A4).")
    p.add_argument("input", help="Input image (PNG/JPG).")
    p.add_argument("--output", "-o", required=True, help="Output base name (without extension).")
    p.add_argument(
        "--format",
        choices=["pptx", "pdf", "both"],
        default="pptx",
        help="Output format (default: pptx).",
    )
    p.add_argument("--dpi", type=int, default=300, help="DPI for panel and output pages (default: 300).")
    p.add_argument("--cols", type=int, default=3, help="Number of columns (default: 3).")
    p.add_argument("--rows", type=int, default=4, help="Number of rows (default: 4).")
    p.add_argument("--margin-mm", type=float, default=10.0, help="Print margins in mm (default: 10).")
    p.add_argument(
        "--shift-x-mm",
        type=float,
        default=0.0,
        help="Horizontal shift of panel content in mm (positive = right).",
    )
    p.add_argument(
        "--shift-y-mm",
        type=float,
        default=0.0,
        help="Vertical shift of panel content in mm (positive = down).",
    )
    p.add_argument(
        "--bg-gray",
        type=str,
        default="#8E8E8E",
        help="Background gray color in #RRGGBB (default: #8E8E8E).",
    )
    p.add_argument(
        "--threshold",
        type=int,
        default=200,
        help="Base threshold for background/foreground separation (0-255).",
    )
    p.add_argument(
        "--detail-delta",
        type=int,
        default=60,
        help="How much below --threshold to cut inner details (default: 60).",
    )
    p.add_argument(
        "--blur",
        type=float,
        default=0.6,
        help="Gaussian blur radius before thresholding (default: 0.6).",
    )
    p.add_argument(
        "--dilate-px",
        type=int,
        default=1,
        help="Morphological closing radius in pixels (default: 1).",
    )
    p.add_argument(
        "--detail-join-px",
        type=int,
        default=2,
        help="Extra closing radius to bridge dotted interior details (default: 2).",
    )
    p.add_argument(
        "--antialias-radius",
        type=float,
        default=0.8,
        help="Base edge smoothing radius in pixels (scaled with output).",
    )
    p.add_argument(
        "--normalize-dpi",
        type=int,
        default=300,
        help="Target DPI for source normalization (0 disables DPI-based scaling).",
    )
    p.add_argument(
        "--normalize-scale",
        type=float,
        default=2.0,
        help="Minimum upscale factor for smoothing jagged edges (default: 2.0).",
    )
    p.add_argument(
        "--normalize-blur",
        type=float,
        default=0.8,
        help=(
            "Baseline Gaussian blur radius used during normalization before adaptive scaling "
            "(default: 0.8). Combine with --normalize-preset noblur to disable smoothing."
        ),
    )
    p.add_argument(
        "--normalize-preset",
        choices=["default", "noblur"],
        default="default",
        help=(
            "Normalization preset. 'default' applies adaptive smoothing; 'noblur' forces "
            "zero blur regardless of other settings."
        ),
    )
    p.add_argument(
        "--fit-mode",
        choices=["fit", "fill", "stretch"],
        default="fit",
        help="How to fit silhouette into panel (default: fit).",
    )
    p.add_argument(
        "--debug-panel",
        action="store_true",
        help="Save debug panel PNG (<output>_panel.png) before tiling.",
    )
    p.add_argument(
        "--invert-mode",
        choices=["auto", "keep", "flip"],
        default="auto",
        help="How to treat silhouette/background inversion (default: auto).",
    )
    p.add_argument(
        "--flip-silhouette",
        action="store_true",
        help="Deprecated alias for --invert-mode=flip.",
    )
    p.add_argument(
        "--verbose",
        action="store_true",
        help="Enable verbose logging (progress information).",
    )
    return p.parse_args(argv)


def parse_bg_color(s: str) -> Tuple[int, int, int]:
    s = s.strip()
    if s.startswith("#"):
        s = s[1:]
    if len(s) != 6:
        raise ValueError(f"Invalid bg-gray value: {s!r}")
    r = int(s[0:2], 16)
    g = int(s[2:4], 16)
    b = int(s[4:6], 16)
    return r, g, b


def main(argv=None) -> None:
    args = parse_args(argv)

    root_logger = logging.getLogger()
    if not root_logger.handlers:
        logging.basicConfig(
            level=logging.INFO,
            format="[%(levelname)s] %(message)s",
            stream=sys.stdout,
        )
    if args.verbose:
        root_logger.setLevel(logging.DEBUG)
    else:
        root_logger.setLevel(logging.INFO)

    inp = Path(args.input)
    if not inp.is_file():
        raise SystemExit(f"Input not found: {inp}")

    out_base = Path(args.output)
    bg_rgb = parse_bg_color(args.bg_gray)

    layout = PanelLayout(
        dpi=args.dpi,
        cols=args.cols,
        rows=args.rows,
        margin_mm=args.margin_mm,
    )

    normalize_blur = max(0.0, float(args.normalize_blur))
    if args.normalize_preset == "noblur":
        normalize_blur = 0.0

    norm_params = NormalizationParams(
        target_dpi=max(0, int(args.normalize_dpi)),
        upscale_factor=max(1.0, float(args.normalize_scale)),
        blur_radius=normalize_blur,
    )

    if args.flip_silhouette:
        invert_mode = "flip"
    else:
        invert_mode = args.invert_mode

    threshold_bg = max(0, min(255, int(args.threshold)))
    threshold_detail = max(0, threshold_bg - int(args.detail_delta))

    bw_params = BWParams(
        threshold_bg=threshold_bg,
        threshold_detail=threshold_detail,
        blur=float(args.blur),
        dilate_px=max(0, int(args.dilate_px)),
        antialias_radius=max(0.0, float(args.antialias_radius)),
        detail_join_px=max(0, int(args.detail_join_px)),
    )

    logger.info("Loading image: %s", inp)
    src_loaded = Image.open(inp).convert("RGB")
    src, norm_scale, norm_blur = normalize_source_image(src_loaded, norm_params)
    if norm_scale > 1.0 or norm_blur > 0:
        logger.info(
            "Source normalization applied (scale %.2f, blur %.2f)",
            norm_scale,
            norm_blur,
        )

    panel, original_panel, auto_flipped = build_white_on_gray_panel(
        src,
        layout=layout,
        bw_params=bw_params,
        bg_rgb=bg_rgb,
        fit_mode=args.fit_mode,
        shift_x_mm=args.shift_x_mm,
        shift_y_mm=args.shift_y_mm,
        invert_mode=invert_mode,
    )

    original_png_name = out_base.name + "_original.png"
    original_path = inp.with_name(original_png_name)
    logger.info("Saving original-size PNG to %s", original_path)
    original_panel.save(original_path, format="PNG", dpi=(layout.dpi, layout.dpi))
    print(f"Original-size PNG: {original_path}")

    if args.debug_panel:
        panel_path = out_base.with_name(out_base.name + "_panel.png")
        panel.save(panel_path, format="PNG", dpi=(layout.dpi, layout.dpi))
        print(f"Debug panel saved to {panel_path}")

    logger.info("Slicing panel into %d tiles", layout.cols * layout.rows)
    tiles = slice_panel_into_tiles(panel, layout)

    if args.format in ("pdf", "both"):
        pdf_path = out_base.with_name(out_base.name + "_3x4_A4.pdf")
        logger.info("Writing PDF to %s", pdf_path)
        save_pdf_from_tiles(tiles, pdf_path, layout, bg_rgb)
        print(f"PDF:  {pdf_path}")

    if args.format in ("pptx", "both"):
        pptx_path = out_base.with_name(out_base.name + "_3x4_A4.pptx")
        logger.info("Writing PPTX to %s", pptx_path)
        save_pptx_from_tiles(tiles, pptx_path, layout, bg_rgb)
        print(f"PPTX: {pptx_path}")

    px_2mm = mm_to_px(2.0, layout.dpi)
    px_14mm = mm_to_px(1.4, layout.dpi)

    if invert_mode == "auto":
        invert_info = "auto->flip" if auto_flipped else "auto->keep"
    else:
        invert_info = invert_mode

    print(
        f"DPI={layout.dpi}, margin={layout.margin_mm} mm, shift={args.shift_x_mm}x{args.shift_y_mm} mm, "
        f"threshold_bg={threshold_bg}, threshold_detail={threshold_detail}, blur={bw_params.blur}, "
        f"dilate={bw_params.dilate_px}px, invert={invert_info}",
    )
    print(f"Rule of thumb: min bridge ~2.0 mm (~{px_2mm}px); hard min ~1.4 mm (~{px_14mm}px).")
    print("Adjust --shift-x-mm/--shift-y-mm if a seam hits a thin detail.")


if __name__ == "__main__":
    try:
        main()
    except Exception:  # pragma: no cover - top-level CLI guard
        logger.exception("Papercut panel builder failed")
        raise
