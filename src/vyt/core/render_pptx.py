"""PPTX rendering stub."""

from __future__ import annotations

from pathlib import Path
from typing import Any, Dict, Iterable

from pptx import Presentation
from pptx.util import Cm


def build(tiles_svg: Iterable[Path], cfg: Dict[str, Any], work: Path, out: Path) -> Path:
    """Generate a minimal PPTX placeholder with one slide per tile."""
    presentation = Presentation()
    presentation.slide_width = Cm(21.0)
    presentation.slide_height = Cm(29.7)

    for index, _ in enumerate(tiles_svg, start=1):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        textbox = slide.shapes.add_textbox(Cm(2), Cm(2), Cm(17), Cm(3))
        textbox.text = f"Tile {index:02d} — EMF/PNG will be here"

    pptx_path = out / f"VYT_{cfg['id']}_A4_pages.pptx"
    presentation.save(pptx_path)
    return pptx_path
